from typing import List, Dict
import logging as log
import os.path
import imghdr
import json
import io

from main import exit_on_failure

import pptx
from pptx.presentation import Presentation
from pptx.util import Cm
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

import numpy as np
import matplotlib.pyplot as plt



class ConfingDecoder:
    """
    ConfigDecoder takes in a ´str´ that is a path to a json file that describes the layout of a presentation.
    """
    def __init__(self, config_path: str):
        self.config_path = config_path
        try:
            json_file = open(self.config_path)
            self.config: dict = json.load(json_file)
        except OSError as err:
            # Error when opening config.json
            log.critical(f"Could not open configuration file, error: {err}")
            exit_on_failure()
        except ValueError as err:
            # Error when parsing config.json
            log.critical(f"The provided file is not in a valid json format, error: {err}")
            exit_on_failure()

        # Json standards usually require keys to be unique. In other cases it appears to behave like, if only the
        # last part is there. So any duplicate key (pres. name) problem is avoided.
        for presentation in self.config.keys():
            # Presentation raises ValueError if the given file is not a pptx file,
            # since we are not providing any, it should be fine
            log.info(f"Creating {presentation}.pptx")
            prs: Presentation = pptx.Presentation()
            for slide in self.config[presentation]:
                match slide["type"]:
                    case "title":
                        log.info(f"Title slide added to {presentation}.pptx")
                        self.__title_slide(prs, slide["title"], slide["content"])
                    case "text":
                        log.info(f"Text slide added to {presentation}.pptx")
                        self.__text_slide(prs, slide["title"], slide["content"])
                    case "list":
                        log.info(f"List slide added to {presentation}.pptx")
                        self.__list_slide(prs, slide["title"], slide["content"])
                    case "picture":
                        img_path = slide["content"]
                        if os.path.isfile(img_path) and imghdr.what(img_path) is not None:
                            log.info(f"Picture slide added to {presentation}.pptx")
                            self.__picture_slide(prs, slide["title"], img_path)
                        else:
                            log.error(f"Content in \"{presentation}\" at {slide}")
                            log.error("The given file path does not exist or unreachable, or the file is not a valid picture by \"imghdr\"")
                            log.info("Skipping slide..")
                    case "plot":
                        csv_path = slide["content"]
                        if os.path.isfile(csv_path):
                            log.info(f"Plot slide added to {presentation}.pptx")
                            self.__plot_slide(prs, slide["title"], csv_path, slide["configuration"])
                        else:
                            log.error("The given file path does not exist or unreachable, or the file is not a valid picture by \"imghdr\"")
                            log.info("Skipping slide..")
                    case other:
                        log.error(f"Unsupported type: \"{other}\" in \"{presentation}\" at {slide}")
                        log.info("Skipping slide..")
                        continue

            prs.save(f"{presentation}.pptx")

    @staticmethod
    def __title_slide(prs, title_name: str, subtitle: str):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title_name
        slide.placeholders[1].text = subtitle

    @staticmethod
    def __text_slide(prs: Presentation, title_name: str, content: str):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title_name

        # LibreOffice Impress seems to have a different kind of coordinate system,
        # so these values are fine-tuned for MS Office only
        w = Cm(22.86)
        h = Cm(12.57)
        x = Cm(1.3)
        y = Cm(4.45)
        textbox = slide.shapes.add_textbox(x, y, w, h)
        tf = textbox.text_frame
        tf.text = content
        tf.add_paragraph()
        tf.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    @staticmethod
    def __list_slide(prs: Presentation, title_name: str, content: List[Dict]):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title_name

        tf = slide.shapes.placeholders[1].text_frame
        for item in content:
            (key, value) = item.items()
            p = tf.add_paragraph()
            p.text = value[1]
            p.level = key[1]

    @staticmethod
    def __picture_slide(prs: Presentation, title_name: str, img_path: str):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title_name

        x = Cm(1.3)
        y = Cm(4.45)

        # Accounting for the possibility of a larger image
        slide.shapes.add_picture(img_path, x, y, prs.slide_height - Cm(10))

        # This one does not, maybe its good maybe not. In my testing it shrinks the not oversized images but not my much.
        # So i guess it's fine
        # slide.shapes.add_picture(img_path, x, y)

    @staticmethod
    def __plot_slide(prs: Presentation, title_name: str, csv_path: str, config: dict):
        arr = np.loadtxt(csv_path, delimiter=";", dtype=float)
        x = []
        y = []
        for i in arr:
            x.append(i[0])
            y.append(i[1])
        plt.plot(x, y)
        plt.ylabel(config["x-label"])
        plt.xlabel(config["y-label"])

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title_name

        x = Cm(1.3)
        y = Cm(4.45)

        image_stream = io.BytesIO()
        plt.savefig(image_stream)
        # slide.shapes.add_picture(image_stream, x, y, prs.slide_height - Cm(10))
        slide.shapes.add_picture(image_stream, x, y)